import os
import io
from typing import Tuple


def parse_file(file_bytes: bytes, filename: str, content_type: str) -> Tuple[str, str, bytes]:
    """
    Parse any uploaded file and return:
    (text_content, file_category, media_bytes)
    file_category: 'text' | 'image' | 'video' | 'unknown'
    """
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
    ct = content_type.lower() if content_type else ''

    # ── VIDEO ────────────────────────────────────────────────────────────
    if ext in ('mp4', 'mov', 'webm') or 'video' in ct:
        return ("", "video", file_bytes)

    # ── IMAGE ────────────────────────────────────────────────────────────
    if ext in ('jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'tiff', 'tif') or 'image' in ct:
        return ("", "image", file_bytes)

    # ── PDF ──────────────────────────────────────────────────────────────
    if ext == 'pdf' or 'pdf' in ct:
        try:
            try:
                import pymupdf as fitz  # newer package name
            except ImportError:
                import fitz              # legacy package name
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return (text, "text", b"")
        except Exception as e:
            return (f"[PDF parse error: {e}]", "text", b"")

    # ── DOCX ─────────────────────────────────────────────────────────────
    if ext == 'docx' or 'wordprocessingml' in ct:
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            return (text, "text", b"")
        except Exception as e:
            return (f"[DOCX parse error: {e}]", "text", b"")

    # ── PPTX ─────────────────────────────────────────────────────────────
    if ext == 'pptx' or 'presentationml' in ct:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        parts.append(shape.text)
            return ('\n'.join(parts), "text", b"")
        except Exception as e:
            return (f"[PPTX parse error: {e}]", "text", b"")

    # ── PLAIN TEXT / CODE ─────────────────────────────────────────────────
    text_exts = {
        'txt', 'md', 'py', 'js', 'ts', 'jsx', 'tsx', 'html', 'css', 'scss',
        'java', 'cpp', 'c', 'h', 'hpp', 'cs', 'go', 'rs', 'rb', 'php', 'swift',
        'kt', 'sql', 'sh', 'bat', 'json', 'xml', 'yaml', 'yml', 'toml', 'ini',
        'csv', 'log', 'tex', 'r', 'dart', 'lua', 'pl', 'scala'
    }
    if ext in text_exts or 'text' in ct:
        try:
            text = file_bytes.decode('utf-8', errors='replace')
            return (text, "text", b"")
        except Exception as e:
            return (f"[Text parse error: {e}]", "text", b"")

    # ── FALLBACK: try UTF-8 decode ────────────────────────────────────────
    try:
        text = file_bytes.decode('utf-8', errors='replace')
        if len(text.strip()) > 50:
            return (text, "text", b"")
    except Exception:
        pass

    return ("", "unknown", b"")


def get_file_icon(ext: str) -> str:
    icons = {
        'pdf': '📄', 'docx': '📝', 'pptx': '📊', 'txt': '📃', 'md': '📋',
        'py': '🐍', 'js': '⚡', 'ts': '🔷', 'html': '🌐', 'css': '🎨',
        'java': '☕', 'cpp': '⚙️', 'c': '⚙️', 'go': '🔵', 'rs': '🦀',
        'jpg': '🖼️', 'jpeg': '🖼️', 'png': '🖼️', 'gif': '🎞️', 'webp': '🖼️',
        'mp4': '🎬', 'mov': '🎬', 'webm': '🎬'
    }
    return icons.get(ext.lower(), '📁')
